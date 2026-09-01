import asyncio
import io
import logging
import os
import random
import re
import time
from collections import defaultdict, deque
from pathlib import Path

import docx
import openpyxl
import requests
from aiogram import Bot, Dispatcher, F
from aiogram.client.default import DefaultBotProperties
from aiogram.enums import ChatAction, ParseMode
from aiogram.filters import Command
from aiogram.types import Message, ReactionTypeEmoji
from aiohttp import web
from google import genai
from google.genai import errors, types
from pptx import Presentation

from config import (
    BOT_TOKEN,
    GEMINI_API_KEY,
    GEMINI_MODEL,
    TAVILY_API_KEY,
    TRIGGER_NAMES,
    SYSTEM_PROMPT,
    HISTORY_SIZE,
    USER_CONTEXT,
)

logging.basicConfig(level=logging.INFO)
logging.getLogger("google_genai.models").setLevel(logging.ERROR)
log = logging.getLogger("Bashma4ek_Bot")

bot = Bot(token=BOT_TOKEN, default=DefaultBotProperties(parse_mode=ParseMode.HTML))
dp = Dispatcher()
ai_client = genai.Client(api_key=GEMINI_API_KEY)


@dp.message.outer_middleware()
async def track_activity_middleware(handler, message: Message, data: dict):
    """Фіксує будь-яке повідомлення від людини в чаті — потрібно окремо від
    основних хендлерів, щоб не забути оновити активність в якомусь з них."""
    if message.chat.type != "private":
        last_human_activity[message.chat.id] = time.time()
        idle_message_sent[message.chat.id] = False
    return await handler(message, data)

# Заповнюється один раз у main() під час старту — щоб не смикати get_me()
# на кожне повідомлення.
BOT_ID: int | None = None
BOT_USERNAME: str = ""
BOT_FULL_NAME: str = ""

START_TIME = time.time()

# Тригер-слова для автоматичного пошуку в інтернеті. Рішення "шукати чи ні"
# приймається в коді напряму (детерміновано), а не моделлю через function
# calling — той підхід виявився нестабільним з gemini-3.1-flash-lite.
SEARCH_TRIGGERS = {
    # укр
    "знайди", "гугл", "пошукай", "новини", "погода", "прогноз",
    "курс", "курси", "долар", "євро", "зарплат", "ціна", "ціни",
    "сьогодні", "зараз", "актуальн", "останні", "свіж",
    "інтернет", "найди", "хто такий", "хто така",
    "коли", "скільки коштує", "де знаходиться", "що сталося", "що відбулось",
    # рос (Влад і Саша частіше пишуть/отримують відповіді російською)
    "найди", "погугли", "поищи", "новост", "прогноз погоды",
    "курс", "доллар", "евро", "зарплат", "цена", "цены",
    "сегодня", "сейчас", "актуальн", "последние", "свеж",
    "интернет", "кто такой", "кто такая",
    "когда", "сколько стоит", "где находится", "что случилось", "что произошло",
    # універсальні / інші мови
    "search", "google",
}

MAX_SEARCH_RESULTS = 5
MAX_FETCH_CHARS = 6000
MAX_DOC_CHARS = 40000

# --- Кеш пошукових запитів ------------------------------------------------
# Якщо кілька людей підряд запитують те саме (наприклад "яка погода?"),
# не варто бити по зовнішньому API двічі — тримаємо результат кілька
# хвилин в пам'яті.
SEARCH_CACHE_TTL_SEC = 5 * 60
_search_cache: dict[str, tuple[float, str]] = {}


def _cache_get(key: str) -> str | None:
    entry = _search_cache.get(key)
    if not entry:
        return None
    ts, value = entry
    if time.time() - ts > SEARCH_CACHE_TTL_SEC:
        _search_cache.pop(key, None)
        return None
    return value


def _cache_set(key: str, value: str) -> None:
    _search_cache[key] = (time.time(), value)
    # Проста самоочистка, щоб словник не ріс нескінченно в довгоживучому процесі.
    if len(_search_cache) > 200:
        oldest_key = min(_search_cache, key=lambda k: _search_cache[k][0])
        _search_cache.pop(oldest_key, None)


# --- Курс валют (НБП — Народний банк Польщі) ------------------------------
# Бот у Польщі, тож база — PLN. Швидше й точніше за пошук по інтернету для
# цієї конкретної, дуже частої категорії запитів.
CURRENCY_TRIGGER_WORDS = {
    "курс", "курси", "курсы", "долар", "доллар", "євро", "евро",
    "гривня", "гривны", "гривень", "злотий", "злотых", "злотого", "фунт",
}

CURRENCY_KEYWORDS = {
    "USD": ["долар", "доллар", "usd", "$"],
    "EUR": ["євро", "евро", "eur", "€"],
    "UAH": ["гривня", "гривны", "гривень", "грн", "uah", "₴"],
    "GBP": ["фунт", "gbp", "£"],
}


def is_currency_query(text: str) -> bool:
    text_lower = text.lower()
    return any(w in text_lower for w in CURRENCY_TRIGGER_WORDS)


def detect_currency_codes(text: str) -> list[str]:
    text_lower = text.lower()
    found = [code for code, keywords in CURRENCY_KEYWORDS.items()
             if any(kw in text_lower for kw in keywords)]
    return found


def _currency_sync(codes: list[str]) -> str:
    codes = codes or ["USD", "EUR"]
    lines = []
    for code in codes:
        try:
            resp = requests.get(
                f"https://api.nbp.pl/api/exchangerates/rates/A/{code}/?format=json",
                timeout=6,
            )
            resp.raise_for_status()
            data = resp.json()
            rate = data["rates"][0]["mid"]
            date = data["rates"][0]["effectiveDate"]
            lines.append(f"- 1 {code} = {rate} PLN (курс НБП станом на {date})")
        except Exception as e:
            log.error(f"Помилка отримання курсу {code} з НБП: {e}")

    if not lines:
        return ""
    return "\n\n[Актуальний курс валют]:\n" + "\n".join(lines)


# --- Погода (Open-Meteo) ---------------------------------------------------
# Безкоштовний API без ключа, точніший і швидший за скрейпінг для цієї
# категорії запитів.
WEATHER_TRIGGER_WORDS = {"погода", "погоду", "погоди", "прогноз погоды", "прогноз погоди"}
DEFAULT_WEATHER_CITY = os.getenv("DEFAULT_WEATHER_CITY", "Bytom")
_WEATHER_STOPWORDS = {
    "сьогодні", "зараз", "завтра", "яка", "буде", "у",
    "сегодня", "сейчас", "завтра", "какая", "будет",
}


def is_weather_query(text: str) -> bool:
    text_lower = text.lower()
    return any(w in text_lower for w in WEATHER_TRIGGER_WORDS)


def extract_weather_city(text: str) -> str:
    """Намагається витягти назву міста після слова 'погода' (напр. 'погода
    у Варшаві'). Якщо не вдалось — використовує місто за замовчуванням."""
    match = re.search(
        r"погод[аиу]?\s*(?:в|у|на)?\s*([A-Za-zА-Яа-яЇїІіЄєҐґ\-]{3,30})",
        text,
        re.IGNORECASE,
    )
    if match:
        candidate = match.group(1).strip()
        if candidate.lower() not in _WEATHER_STOPWORDS:
            return candidate
    return DEFAULT_WEATHER_CITY


def _weather_sync(city: str) -> str:
    try:
        geo_resp = requests.get(
            "https://geocoding-api.open-meteo.com/v1/search",
            params={"name": city, "count": 1, "language": "uk"},
            timeout=6,
        )
        geo_resp.raise_for_status()
        results = geo_resp.json().get("results")
        if not results:
            return ""
        loc = results[0]
        found_name = loc.get("name", city)
        country = loc.get("country", "")

        forecast_resp = requests.get(
            "https://api.open-meteo.com/v1/forecast",
            params={
                "latitude": loc["latitude"],
                "longitude": loc["longitude"],
                "current": "temperature_2m,apparent_temperature,precipitation,wind_speed_10m",
                "timezone": "auto",
            },
            timeout=6,
        )
        forecast_resp.raise_for_status()
        current = forecast_resp.json().get("current", {})
        if not current:
            return ""

        return (
            f"\n\n[Поточна погода — {found_name}, {country}]:\n"
            f"Температура: {current.get('temperature_2m')}°C "
            f"(відчувається як {current.get('apparent_temperature')}°C)\n"
            f"Вітер: {current.get('wind_speed_10m')} км/год, "
            f"опади: {current.get('precipitation')} мм"
        )
    except Exception as e:
        log.error(f"Помилка отримання погоди для {city!r}: {e}")
        return ""


# --- Загальний пошук в інтернеті (Tavily) ----------------------------------
# Tavily заточений під LLM-агентів: одразу повертає очищений релевантний
# контент по кожному результату (не треба окремо парсити HTML сторінки, як
# із сирими сніпетами DuckDuckGo).
def _tavily_search_sync(query: str) -> str:
    if not TAVILY_API_KEY:
        log.error("TAVILY_API_KEY не задано в .env — пошук в інтернеті вимкнено")
        return ""

    try:
        resp = requests.post(
            "https://api.tavily.com/search",
            json={
                "api_key": TAVILY_API_KEY,
                "query": query,
                "search_depth": "basic",
                "max_results": MAX_SEARCH_RESULTS,
                "include_answer": True,
            },
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()
    except Exception as e:
        log.error(f"Помилка пошуку Tavily: {e}")
        return ""

    lines = ["\n\n[Знайдена актуальна інформація з інтернету]:"]

    answer = data.get("answer")
    if answer:
        lines.append(f"Коротка відповідь: {answer}")

    for r in data.get("results", []):
        title = r.get("title", "")
        content = r.get("content", "")
        url = r.get("url", "")
        if len(content) > MAX_FETCH_CHARS:
            content = content[:MAX_FETCH_CHARS] + "...[текст обрізано]"
        lines.append(f"- {title}: {content} ({url})")

    if len(lines) == 1:
        return ""
    return "\n".join(lines)

# Скільки разів повторити запит до Gemini при тимчасових помилках
# (мережа/сервер), перш ніж здатися.
GEMINI_MAX_RETRIES = 2
GEMINI_RETRY_DELAY = 2.0

# Емодзі-реакції, дозволені Telegram Bot API для звичайних (не преміум)
# ботів. Список неповний, але покриває базові емоції — цього достатньо,
# щоб бот міг "просто лайкнути" замість писати текст.
ALLOWED_REACTIONS = {
    "👍", "👎", "❤", "🔥", "🥰", "👏", "😁", "🤔", "🤯", "😱",
    "🤬", "😢", "🎉", "🤩", "🤮", "💩", "🙏", "👌", "🕊", "🤡",
    "🥱", "🥴", "😍", "🐳", "❤‍🔥", "🌚", "🌭", "💯", "🤣", "⚡",
    "🍌", "🏆", "💔", "🤨", "😐", "🍓", "🍾", "💋", "🖕", "😈",
    "😴", "😭", "🤓", "👻", "👨‍💻", "👀", "🎃", "🙈", "😇", "😨",
    "🤝", "✍", "🤗", "🫡", "🎅", "🎄", "☃", "💅", "🤪", "🗿",
    "🆒", "💘", "🙉", "🦄", "😘", "💊", "🙊", "😎", "👾", "🤷‍♂",
    "🤷", "🤷‍♀", "😡",
}

# Маркер, яким модель може позначити "хочу відповісти реакцією, а не
# текстом" — детально в SYSTEM_PROMPT-доповненні в ask_gemini().
REACTION_PREFIX = "REACTION:"

# --- Проактивні повідомлення в тихому чаті ------------------------------
# Якщо в чаті тиша довше IDLE_HOURS годин — бот раз може сам щось написати
# (без згадки), а потім чекає нової активності від людей, перш ніж це
# зможе повторитися знову (щоб не спамити в мертвий чат щогодини).
IDLE_HOURS = float(os.getenv("IDLE_HOURS", "7"))
IDLE_CHECK_INTERVAL_SEC = 15 * 60

# chat_id -> час останнього повідомлення від людини (time.time())
last_human_activity: dict[int, float] = {}
# chat_id -> чи вже "вистрелили" проактивним повідомленням за цей період тиші
idle_message_sent: dict[int, bool] = {}

# --- Реакції на повідомлення без тегу бота -------------------------------
# Незалежно від того, тегнули бота чи ні, він може іноді (не на кожне
# повідомлення) поставити емодзі-реакцію, якщо вважає це доречним.
# Рішення "доречно чи ні" приймає сама модель окремим дешевим запитом
# (max_output_tokens=10), а ймовірність нижче — це фільтр ДО запиту, щоб
# не смикати API на кожну репліку в чаті.
REACT_UNPROMPTED_ENABLED = os.getenv("REACT_UNPROMPTED_ENABLED", "true").lower() == "true"
REACT_UNPROMPTED_CHANCE = float(os.getenv("REACT_UNPROMPTED_CHANCE", "0.35"))


def needs_web_search(text: str) -> bool:
    """Перевіряє, чи варто автоматично зробити пошук в інтернеті."""
    text_lower = (text or "").lower()
    return any(trigger in text_lower for trigger in SEARCH_TRIGGERS)


async def get_web_context(query: str) -> str:
    """Роутер: спочатку перевіряє кеш, далі — чи це запит про курс валют
    або погоду (спеціалізовані швидкі й точні джерела), і лише якщо ні —
    йде в загальний пошук через Tavily. Успішний результат кешується на
    SEARCH_CACHE_TTL_SEC, щоб однаковий запит від різних людей підряд не
    бив по зовнішніх API двічі."""
    cache_key = query.strip().lower()
    cached = _cache_get(cache_key)
    if cached is not None:
        log.info(f"Пошук: кеш-хіт для запиту {query!r}")
        return cached

    result = ""

    if is_currency_query(query):
        codes = detect_currency_codes(query)
        result = await asyncio.to_thread(_currency_sync, codes)

    elif is_weather_query(query):
        city = extract_weather_city(query)
        result = await asyncio.to_thread(_weather_sync, city)

    if not result:
        result = await asyncio.to_thread(_tavily_search_sync, query)

    if result:
        _cache_set(cache_key, result)

    return result


# chat_id -> deque of {"role": ..., "content": ...}
history: dict[int, deque] = defaultdict(lambda: deque(maxlen=HISTORY_SIZE))

_NAMES_ALT = "|".join(re.escape(n) for n in TRIGGER_NAMES)
NAME_PATTERN = re.compile(
    rf"(^\s*({_NAMES_ALT})\b)|(\b({_NAMES_ALT})\s*[,!?.]*\s*$)",
    re.IGNORECASE,
)


def strip_trigger(text: str, bot_username: str) -> str:
    """Прибирає @згадку бота та тригер-ім'я з тексту питання."""
    text = text.replace(f"@{bot_username}", "")
    text = NAME_PATTERN.sub("", text, count=1)
    return text.strip(" ,:.!?-")


def strip_name_prefix(text: str, sender: str, bot_name: str) -> str:
    """Прибирає префікс на кшталт 'koserpa: ' або 'Башмак: ' з відповіді моделі."""
    text = text.strip()
    names = "|".join(re.escape(n) for n in {sender, bot_name, *TRIGGER_NAMES} if n)
    text = re.sub(rf"^\s*(?:{names})\s*:\s*", "", text, count=1, flags=re.IGNORECASE)
    return text.strip()


def parse_reaction_answer(answer: str) -> str | None:
    """Якщо відповідь моделі — це маркер REACTION:<емодзі>, повертає сам
    емодзі (якщо він у дозволеному списку). Інакше None."""
    stripped = answer.strip()
    if not stripped.startswith(REACTION_PREFIX):
        return None
    emoji = stripped[len(REACTION_PREFIX):].strip()
    if emoji in ALLOWED_REACTIONS:
        return emoji
    log.warning(f"Модель попросила недозволену реакцію: {emoji!r}, ігнорую маркер")
    return None


def was_mentioned(message: Message) -> bool:
    text = message.text or message.caption

    if message.reply_to_message and message.reply_to_message.from_user:
        if message.reply_to_message.from_user.id == BOT_ID:
            return True

    if not text:
        return False

    if BOT_USERNAME and f"@{BOT_USERNAME}".lower() in text.lower():
        return True

    if NAME_PATTERN.search(text):
        return True

    return False


async def ask_gemini(contents: list) -> str:
    """Викликає Gemini API. Пошук в інтернеті вже підмішаний у текст промпту
    заздалегідь (детерміновано, у хендлерах) — сюди він приходить готовим.
    При тимчасових (мережа/сервер) помилках робить кілька повторних спроб."""
    last_error: Exception | None = None

    for attempt in range(GEMINI_MAX_RETRIES + 1):
        try:
            response = await ai_client.aio.models.generate_content(
                model=GEMINI_MODEL,
                contents=contents,
                config=types.GenerateContentConfig(
                    system_instruction=SYSTEM_PROMPT
                    + " У переписці повідомлення користувачів позначені як "
                      "'Ім'я: текст' — звертай увагу, хто саме що написав, "
                      "але у своїй відповіді імена дублювати не треба. "
                      ""
                      "ІНОДІ доречніше не писати текст, а просто поставити "
                      "емодзі-реакцію на повідомлення (наприклад, коротке "
                      "'ахах', 'лол', '+', 'згоден', жарт що не вартий "
                      "розгорнутої відповіді, чи щось шокуюче/смішне). "
                      "Якщо вирішив відповісти реакцією — виведи ЛИШЕ рядок "
                      f"'{REACTION_PREFIX}<емодзі>' і нічого більше, без "
                      "жодного тексту до чи після. Дозволені емодзі: "
                      + " ".join(sorted(ALLOWED_REACTIONS))
                      + ". Не зловживай цим — переважно все ж пиши звичайну "
                        "текстову відповідь, реакція лише коли вона реально "
                        "доречніша за слова. "
                      ""
                      "Якщо у повідомленні через @ тегнуто кількох людей "
                      "одразу (і для них є [Про згаданих людей] в промпті) — "
                      "можеш відповісти, врахувавши обох/усіх, а не тільки "
                      "того, хто писав. "
                      ""
                      "Коли тобі надсилають фото, стікер, гіфку, голосове чи "
                      "файл (зокрема через reply на чиєсь повідомлення з "
                      "позначкою '[у відповідь на ...]') — НЕ роби сухий "
                      "'звіт' про вміст. Реагуй на це як жива людина в "
                      "переписці: постьобатися, здивуватись, оцінити, "
                      "прокоментувати по суті — залежно від того, що на "
                      "фото/в файлі і в якому режимі тону зараз розмова. "
                      "Опис вмісту — лише якщо він реально потрібен для "
                      "відповіді, а не сама мета відповіді.",
                    max_output_tokens=600,
                    temperature=0.7,
                ),
            )
            return (response.text or "").strip()

        except errors.ClientError as e:
            if e.code == 429:
                log.warning("Запит відхилено: ліміт 429 (RESOURCE_EXHAUSTED)")
                return "Зараз отримую занадто багато запитів 🤯. Зачекай 1-2 хвилини!"
            log.error(f"Помилка Gemini API (ClientError, без повтору): {e}")
            return "Виникла помилка при зверненні до AI 😔"

        except errors.ServerError as e:
            last_error = e
            log.warning(
                f"Тимчасова помилка Gemini API (спроба {attempt + 1}/"
                f"{GEMINI_MAX_RETRIES + 1}): {e}"
            )
        except Exception as e:
            last_error = e
            log.warning(
                f"Несподівана помилка в ask_gemini (спроба {attempt + 1}/"
                f"{GEMINI_MAX_RETRIES + 1}): {e}"
            )

        if attempt < GEMINI_MAX_RETRIES:
            await asyncio.sleep(GEMINI_RETRY_DELAY)

    log.error(f"ask_gemini: усі спроби вичерпано, остання помилка: {last_error}")
    return "Не вдалося сформулювати відповідь 😔"


async def transcribe_media(data: bytes, mime_type: str, kind_label: str) -> str:
    """Спільна логіка транскрибування аудіо/відео (voice / video_note)."""
    contents = [
        {
            "role": "user",
            "parts": [
                {
                    "text": "Транскрибуй мовлення з цього медіа дослівно, "
                    "тією мовою, якою його промовлено. У відповідь дай ТІЛЬКИ "
                    "текст транскрипції, без жодних коментарів чи лапок."
                },
                types.Part.from_bytes(data=data, mime_type=mime_type),
            ],
        }
    ]
    try:
        response = await ai_client.aio.models.generate_content(
            model=GEMINI_MODEL,
            contents=contents,
            config=types.GenerateContentConfig(max_output_tokens=400, temperature=0.2),
        )
        return (response.text or "").strip()
    except Exception:
        log.exception(f"Транскрибування ({kind_label}) не вдалося")
        return ""


def get_sender_context(message: Message) -> str:
    """Повертає підказку боту про відправника, якщо це відомий учасник."""
    if not message.from_user or not message.from_user.username:
        return ""
    username = message.from_user.username.lstrip("@").lower()
    for known_username, context in USER_CONTEXT.items():
        if known_username.lower() == username:
            return context
    return ""


def get_mentioned_users_context(text: str) -> str:
    """Шукає в тексті @згадки відомих учасників (окрім самого бота) і
    повертає для них контекст із USER_CONTEXT — щоб бот міг врахувати
    кількох людей одразу, а не тільки того, хто написав повідомлення."""
    if not text:
        return ""

    mentioned_usernames = set(re.findall(r"@(\w+)", text))
    bot_username_lower = (BOT_USERNAME or "").lower()

    blocks = []
    for username in mentioned_usernames:
        if username.lower() == bot_username_lower:
            continue
        for known_username, context in USER_CONTEXT.items():
            if known_username.lower() == username.lower():
                blocks.append(f"@{known_username}: {context}")
                break

    if not blocks:
        return ""
    return "\n[Про згаданих людей]:\n" + "\n".join(blocks)


def extract_document_text(file_name: str, data: bytes, mime_type: str | None):
    """Готує вміст файлу для Gemini."""
    ext = Path(file_name or "").suffix.lower()
    mime_type = mime_type or ""

    if ext == ".pdf" or mime_type == "application/pdf":
        return None, types.Part.from_bytes(data=data, mime_type="application/pdf")

    if ext == ".docx" or "wordprocessingml.document" in mime_type:
        document = docx.Document(io.BytesIO(data))
        lines = [p.text for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                lines.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(lines), None

    if ext == ".xlsx" or "spreadsheetml.sheet" in mime_type:
        workbook = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        lines = []
        for sheet in workbook.worksheets:
            lines.append(f"# Аркуш: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    lines.append(" | ".join("" if c is None else str(c) for c in row))
        return "\n".join(lines), None

    if ext == ".pptx" or "presentationml.presentation" in mime_type:
        presentation = Presentation(io.BytesIO(data))
        lines = []
        for i, slide in enumerate(presentation.slides, start=1):
            lines.append(f"# Слайд {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        lines.append(text)
        return "\n".join(lines), None

    text_exts = {".txt", ".md", ".csv", ".json", ".log", ".py", ".yaml", ".yml", ".xml", ".html"}
    if ext in text_exts or mime_type.startswith("text/"):
        return data.decode("utf-8", errors="ignore"), None

    return None, None


async def build_reply_media_context(replied: Message) -> tuple[list, str]:
    """Якщо повідомлення, на яке відповіли (reply), містить фото/стікер/
    гіфку/файл/голосове/кружок — підвантажує цей вміст і повертає
    (extra_parts, опис для промпту). Завдяки цьому 'reply на фото + тег
    бота в тексті' сприймається так само, ніби фото щойно надіслали й
    одразу тегнули бота під ним."""
    extra_parts: list = []
    description = ""

    try:
        if replied.photo:
            photo = replied.photo[-1]
            file = await bot.get_file(photo.file_id)
            buffer = await bot.download_file(file.file_path)
            extra_parts.append(types.Part.from_bytes(data=buffer.read(), mime_type="image/jpeg"))
            description = "[у відповідь на фото]"
            if replied.caption:
                description += f" (підпис до фото: {replied.caption})"

        elif replied.sticker and not (replied.sticker.is_animated or replied.sticker.is_video):
            file = await bot.get_file(replied.sticker.file_id)
            buffer = await bot.download_file(file.file_path)
            extra_parts.append(types.Part.from_bytes(data=buffer.read(), mime_type="image/webp"))
            description = f"[у відповідь на стікер {replied.sticker.emoji or ''}]"

        elif replied.sticker:
            # анімовані/відео-стікери vision не читає — тільки емодзі
            description = f"[у відповідь на анімований стікер {replied.sticker.emoji or ''}]"

        elif replied.animation:
            file = await bot.get_file(replied.animation.file_id)
            buffer = await bot.download_file(file.file_path)
            extra_parts.append(types.Part.from_bytes(data=buffer.read(), mime_type="video/mp4"))
            description = "[у відповідь на гіфку]"
            if replied.caption:
                description += f" (підпис до гіфки: {replied.caption})"

        elif replied.document:
            file_name = replied.document.file_name or "файл"
            file = await bot.get_file(replied.document.file_id)
            buffer = await bot.download_file(file.file_path)
            data = buffer.read()
            text_content, raw_part = extract_document_text(file_name, data, replied.document.mime_type)
            description = f"[у відповідь на файл {file_name}]"
            if raw_part is not None:
                extra_parts.append(raw_part)
            elif text_content:
                trimmed = text_content[:MAX_DOC_CHARS]
                if len(text_content) > MAX_DOC_CHARS:
                    trimmed += "\n...[текст обрізано, файл завеликий]"
                description += f"\nВміст файлу:\n{trimmed}"
            else:
                description += " (формат файлу не підтримується)"

        elif replied.voice:
            file = await bot.get_file(replied.voice.file_id)
            buffer = await bot.download_file(file.file_path)
            transcript = await transcribe_media(buffer.read(), "audio/ogg", "voice-reply")
            if transcript:
                description = f"[у відповідь на голосове]: {transcript}"

        elif replied.video_note:
            file = await bot.get_file(replied.video_note.file_id)
            buffer = await bot.download_file(file.file_path)
            transcript = await transcribe_media(buffer.read(), "video/mp4", "video_note-reply")
            if transcript:
                description = f"[у відповідь на кружок]: {transcript}"

    except Exception:
        log.exception("Не вдалось підвантажити медіа з reply_to_message")
        return [], ""

    return extra_parts, description


async def maybe_react_unprompted(message: Message, sender: str, content_text: str) -> None:
    """Може поставити емодзі-реакцію на будь-яке повідомлення, навіть якщо
    бота не тегали — вибірково і лише коли це реально доречно. Викликається
    як fire-and-forget задача, нічого не пише в чат і не чіпає history."""
    if not REACT_UNPROMPTED_ENABLED:
        return
    if not content_text or not content_text.strip():
        return
    if random.random() > REACT_UNPROMPTED_CHANCE:
        return

    prompt = (
        f"Повідомлення в чаті від {sender}: \"{content_text}\"\n\n"
        "Чи варто відреагувати на нього емодзі-реакцією (без тексту)? "
        "Це доречно ЛИШЕ для дійсно яскравих випадків: дуже смішне, "
        "шокуюче, влучне, драма, класна новина, бʼющий факап тощо. "
        "Для нейтральних, буденних чи незрозумілих повідомлень — реакція "
        "НЕ потрібна, це має бути рідкісна дія, а не звичка. "
        "Якщо доречно — виведи РІВНО ОДНЕ емодзі з цього списку: "
        + " ".join(sorted(ALLOWED_REACTIONS))
        + ". Якщо ні — виведи рівно слово NONE. Без пояснень і зайвих символів."
    )

    try:
        response = await ai_client.aio.models.generate_content(
            model=GEMINI_MODEL,
            contents=[{"role": "user", "parts": [{"text": prompt}]}],
            config=types.GenerateContentConfig(max_output_tokens=10, temperature=0.4),
        )
        answer = (response.text or "").strip()
    except Exception:
        log.exception("Не вдалось перевірити доречність незапитаної реакції")
        return

    if answer not in ALLOWED_REACTIONS:
        return

    try:
        await bot.set_message_reaction(
            chat_id=message.chat.id,
            message_id=message.message_id,
            reaction=[ReactionTypeEmoji(emoji=answer)],
        )
    except Exception:
        log.exception("Не вдалось поставити незапитану реакцію")


# ---------------------------------------------------------------------------
# Спільна логіка для всіх типів повідомлень.
#
# Кожен хендлер (текст/фото/стікер/гіфка/кружок/голосове/документ) робить
# по суті одне й те саме: сформувати текстовий опис події, за потреби додати
# веб-контекст і контекст про відправника, звернутись до Gemini і записати
# результат в історію чату. Різниця лише в тому, як саме будується
# "сирий" контент (question_text) і які додаткові Part-и (фото/відео/файл)
# додаються в запит.
# ---------------------------------------------------------------------------


async def process_and_reply(
    message: Message,
    sender: str,
    question_text: str,
    *,
    extra_parts: list | None = None,
    history_label: str,
) -> None:
    """Формує повний промпт, звертається до Gemini і відповідає в чат.
    Викликається тільки коли бот був згаданий (mentioned=True)."""
    chat_history = history[message.chat.id]

    full_prompt_text = f"{sender}: {question_text}"
    if needs_web_search(question_text):
        web_info = await get_web_context(question_text)
        if web_info:
            full_prompt_text += web_info

    sender_context = get_sender_context(message)
    if sender_context:
        full_prompt_text += f"\n[Про співрозмовника: {sender_context}]"

    mentioned_context = get_mentioned_users_context(question_text)
    if mentioned_context:
        full_prompt_text += mentioned_context

    parts = [{"text": full_prompt_text}]
    if extra_parts:
        parts.extend(extra_parts)

    contents = list(chat_history)
    contents.append({"role": "user", "parts": parts})

    await bot.send_chat_action(message.chat.id, ChatAction.TYPING)

    try:
        answer = (await ask_gemini(contents)) or "Не зміг сформулювати відповідь 🤔"
        answer = strip_name_prefix(answer, sender, BOT_FULL_NAME)
    except Exception:
        log.exception("AI request failed")
        answer = "Вибач, сталася помилка при зверненні до AI 😔"

    chat_history.append(
        {"role": "user", "parts": [{"text": f"{sender}: {history_label}"}]}
    )

    reaction_emoji = parse_reaction_answer(answer)
    if reaction_emoji:
        # В історію кладемо коротку текстову позначку, а не сирий маркер —
        # інакше модель побачить "REACTION:🔥" в наступному контексті й
        # може почати його копіювати як звичайний текст.
        chat_history.append(
            {"role": "model", "parts": [{"text": f"(відреагував {reaction_emoji})"}]}
        )
        try:
            await bot.set_message_reaction(
                chat_id=message.chat.id,
                message_id=message.message_id,
                reaction=[ReactionTypeEmoji(emoji=reaction_emoji)],
            )
        except Exception:
            log.exception("Не вдалось поставити реакцію, відповідаю текстом")
            await message.reply(reaction_emoji)
        return

    chat_history.append({"role": "model", "parts": [{"text": answer}]})
    await message.reply(answer)


def remember_only(message: Message, sender: str, note: str) -> None:
    """Записує подію в історію чату без звернення до Gemini (коли бота не
    згадали). Додатково запускає фонову перевірку — чи не варто все ж
    відреагувати емодзі на це повідомлення (без тегу)."""
    history[message.chat.id].append(
        {"role": "user", "parts": [{"text": f"{sender}: {note}"}]}
    )
    asyncio.create_task(maybe_react_unprompted(message, sender, note))


async def send_idle_message(chat_id: int) -> None:
    """Формує і надсилає одне проактивне повідомлення в тихий чат."""
    chat_history = history[chat_id]
    contents = list(chat_history)
    contents.append(
        {
            "role": "user",
            "parts": [{
                "text": (
                    f"[СИСТЕМНЕ]: у чаті тиша вже {IDLE_HOURS:.0f}+ годин. "
                    "Напиши щось одне коротке від себе, щоб оживити чат — "
                    "жарт, провокаційне питання чи коротку думку, за темою "
                    "останніх повідомлень якщо вони були, у своєму "
                    "звичному стилі. Рівно 1 речення. Без звернення до "
                    "когось конкретного і без пояснень, що ти бот, що чат "
                    "мовчав, чи щось подібне — просто природне повідомлення."
                )
            }],
        }
    )

    answer = await ask_gemini(contents)
    if not answer or parse_reaction_answer(answer):
        # Реакцію тут ставити нема на що (це не відповідь комусь конкретно),
        # тож у цьому випадку просто пропускаємо тик.
        return

    answer = strip_name_prefix(answer, "", BOT_FULL_NAME)
    chat_history.append({"role": "model", "parts": [{"text": answer}]})
    await bot.send_message(chat_id, answer)


async def idle_chat_watcher():
    """Раз на IDLE_CHECK_INTERVAL_SEC проходиться по відомих чатах: якщо
    тиша довша за IDLE_HOURS і бот ще не писав за цей період тиші —
    надсилає одне проактивне повідомлення."""
    while True:
        await asyncio.sleep(IDLE_CHECK_INTERVAL_SEC)
        now = time.time()

        for chat_id, last_active in list(last_human_activity.items()):
            if idle_message_sent.get(chat_id):
                continue
            if now - last_active < IDLE_HOURS * 3600:
                continue

            # Ставимо прапорець одразу (до відправки), щоб не задвоїти
            # повідомлення, якщо цей тик з якоїсь причини затягнеться.
            idle_message_sent[chat_id] = True
            try:
                await send_idle_message(chat_id)
            except Exception:
                log.exception(f"Не вдалось надіслати проактивне повідомлення в чат {chat_id}")


@dp.message(Command("start", "help"))
async def cmd_start(message: Message):
    await message.answer(
        "Привіт! Я бот-асистент. Я запам'ятовую переписку в чаті, "
        "а відповідаю, коли мене тегнуть (@бот) або звертаються по імені "
        f"({', '.join(TRIGGER_NAMES)})."
    )


@dp.message(Command("reset"))
async def cmd_reset(message: Message):
    history[message.chat.id].clear()
    await message.answer("Пам'ять цього чату очищена 🧹")


@dp.message(Command("status"))
async def cmd_status(message: Message):
    uptime_sec = int(time.time() - START_TIME)
    hours, remainder = divmod(uptime_sec, 3600)
    minutes, seconds = divmod(remainder, 60)
    chat_len = len(history[message.chat.id])

    await message.answer(
        "📊 Статус бота\n"
        f"Модель: {GEMINI_MODEL}\n"
        f"Uptime: {hours}г {minutes}хв {seconds}с\n"
        f"Повідомлень в пам'яті цього чату: {chat_len}/{HISTORY_SIZE}"
    )


@dp.message(F.text)
async def handle_message(message: Message):
    sender = message.from_user.full_name if message.from_user else "Хтось"
    mentioned = was_mentioned(message)

    if not mentioned:
        remember_only(message, sender, message.text)
        return

    question = strip_trigger(message.text, BOT_USERNAME)
    if not question:
        question = "Привіт! Про що поговоримо?"

    # Якщо це reply на чиєсь фото/стікер/гіфку/файл/голосове/кружок і в
    # тексті тегнули бота — підвантажуємо той медіа-контент, щоб бот
    # сприйняв це так, ніби йому щойно надіслали й одразу тегнули.
    extra_parts: list = []
    replied = message.reply_to_message
    if replied and (not replied.from_user or replied.from_user.id != BOT_ID):
        reply_parts, reply_description = await build_reply_media_context(replied)
        if reply_description:
            question = f"{question}\n{reply_description}" if question else reply_description
            extra_parts = reply_parts

    await process_and_reply(
        message, sender, question,
        extra_parts=extra_parts or None,
        history_label=question,
    )


@dp.message(F.photo)
async def handle_photo(message: Message):
    sender = message.from_user.full_name if message.from_user else "Хтось"
    caption = message.caption or ""
    mentioned = was_mentioned(message)

    if not mentioned:
        note = f"[надіслав(-ла) фото]"
        if caption:
            note += f" {caption}"
        remember_only(message, sender, note)
        return

    question = strip_trigger(caption, BOT_USERNAME) or "Що на цьому фото?"

    try:
        photo = message.photo[-1]
        file = await bot.get_file(photo.file_id)
        buffer = await bot.download_file(file.file_path)
        image_bytes = buffer.read()
    except Exception:
        log.exception("Failed to download photo")
        await message.reply("Не вдалось завантажити фото 😔")
        return

    extra_parts = [types.Part.from_bytes(data=image_bytes, mime_type="image/jpeg")]
    await process_and_reply(
        message, sender, question,
        extra_parts=extra_parts,
        history_label=f"[фото] {question}",
    )


@dp.message(F.sticker)
async def handle_sticker(message: Message):
    sender = message.from_user.full_name if message.from_user else "Хтось"
    sticker = message.sticker
    emoji = sticker.emoji or "🙂"
    mentioned = was_mentioned(message)

    # анімовані (.tgs) і відео-стікери (.webm) Gemini vision напряму не їсть —
    # фіксуємо тільки емодзі, без реального аналізу картинки
    if sticker.is_animated or sticker.is_video:
        if not mentioned:
            remember_only(message, sender, f"[анімований стікер {emoji}]")
            return
        await process_and_reply(
            message, sender, f"[надіслав(-ла) анімований стікер {emoji}]",
            history_label=f"[анімований стікер {emoji}]",
        )
        return

    # звичайний статичний стікер — webp, Gemini vision їсть напряму
    try:
        file = await bot.get_file(sticker.file_id)
        buffer = await bot.download_file(file.file_path)
        sticker_bytes = buffer.read()
    except Exception:
        log.exception("Failed to download sticker")
        if mentioned:
            await message.reply("Не вдалось завантажити стікер 😔")
        return

    if not mentioned:
        remember_only(message, sender, f"[надіслав(-ла) стікер {emoji}]")
        return

    extra_parts = [types.Part.from_bytes(data=sticker_bytes, mime_type="image/webp")]
    await process_and_reply(
        message, sender, f"[надіслав(-ла) стікер, емодзі: {emoji}]",
        extra_parts=extra_parts,
        history_label=f"[надіслав(-ла) стікер {emoji}]",
    )


@dp.message(F.animation)
async def handle_animation(message: Message):
    """GIF в Telegram технічно приходить як mp4 без звуку (F.animation) —
    gemini-3.1-flash-lite підтримує відео на вході, тож кидаємо файл напряму,
    без потреби витягувати кадр через ffmpeg."""
    sender = message.from_user.full_name if message.from_user else "Хтось"
    caption = message.caption or ""
    mentioned = was_mentioned(message)

    if not mentioned:
        note = "[надіслав(-ла) гіфку]"
        if caption:
            note += f" {caption}"
        remember_only(message, sender, note)
        return

    animation = message.animation
    if animation.file_size and animation.file_size > 20 * 1024 * 1024:
        await message.reply("Гіфка більша за 20 МБ — стільки бот завантажити не може 😔")
        return

    try:
        file = await bot.get_file(animation.file_id)
        buffer = await bot.download_file(file.file_path)
        animation_bytes = buffer.read()
    except Exception:
        log.exception("Failed to download animation")
        await message.reply("Не вдалось завантажити гіфку 😔")
        return

    question = strip_trigger(caption, BOT_USERNAME) or "Що відбувається на цій гіфці?"
    extra_parts = [types.Part.from_bytes(data=animation_bytes, mime_type="video/mp4")]
    await process_and_reply(
        message, sender, question,
        extra_parts=extra_parts,
        history_label=f"[гіфка] {question}",
    )


@dp.message(F.video_note)
async def handle_video_note(message: Message):
    """Кружки (video_note): транскрибуємо мовлення так само, як голосові —
    в історію і у відповідь йде тільки текст транскрипції, без самого відео."""
    sender = message.from_user.full_name if message.from_user else "Хтось"
    mentioned = was_mentioned(message)

    try:
        video_note = message.video_note
        file = await bot.get_file(video_note.file_id)
        buffer = await bot.download_file(file.file_path)
        video_note_bytes = buffer.read()
    except Exception:
        log.exception("Failed to download video note")
        if mentioned:
            await message.reply("Не вдалось завантажити кружок 😔")
        return

    transcript = await transcribe_media(video_note_bytes, "video/mp4", "video_note")

    if not transcript:
        if mentioned:
            await message.reply("Не вдалось розпізнати кружок 😔")
        else:
            remember_only(message, sender, "[кружок]")
        return

    if not mentioned:
        remember_only(message, sender, f"[кружок] {transcript}")
        return

    await process_and_reply(
        message, sender, transcript,
        history_label=f"[кружок] {transcript}",
    )


@dp.message(F.voice)
async def handle_voice(message: Message):
    """Голосові повідомлення: Gemini сам транскрибує та розуміє аудіо напряму."""
    sender = message.from_user.full_name if message.from_user else "Хтось"
    caption = message.caption or ""
    mentioned = was_mentioned(message)

    try:
        voice = message.voice
        file = await bot.get_file(voice.file_id)
        buffer = await bot.download_file(file.file_path)
        voice_bytes = buffer.read()
    except Exception:
        log.exception("Failed to download voice message")
        if mentioned:
            await message.reply("Не вдалось завантажити голосове 😔")
        return

    transcript = await transcribe_media(voice_bytes, "audio/ogg", "voice")

    if not transcript:
        if mentioned:
            await message.reply("Не вдалось розпізнати голосове повідомлення 😔")
        else:
            remember_only(message, sender, "[голосове повідомлення]")
        return

    if not mentioned:
        note = f"[голосове] {transcript}"
        if caption:
            note += f" ({caption})"
        remember_only(message, sender, note)
        return

    question = strip_trigger(caption, BOT_USERNAME)
    question_text = transcript
    if question:
        question_text += f"\n[Коментар до голосового]: {question}"

    await process_and_reply(
        message, sender, question_text,
        history_label=f"[голосове] {transcript}",
    )


@dp.message(F.document)
async def handle_document(message: Message):
    sender = message.from_user.full_name if message.from_user else "Хтось"
    caption = message.caption or ""
    doc = message.document
    file_name = doc.file_name or "файл"
    mentioned = was_mentioned(message)

    if not mentioned:
        note = f"[надіслав(-ла) файл {file_name}]"
        if caption:
            note += f" {caption}"
        remember_only(message, sender, note)
        return

    if doc.file_size and doc.file_size > 20 * 1024 * 1024:
        await message.reply("Файл більший за 20 МБ — стільки бот завантажити не може 😔")
        return

    try:
        tg_file = await bot.get_file(doc.file_id)
        buffer = await bot.download_file(tg_file.file_path)
        data = buffer.read()
    except Exception:
        log.exception("Failed to download document")
        await message.reply("Не вдалось завантажити файл 😔")
        return

    text_content, raw_part = extract_document_text(file_name, data, doc.mime_type)

    if text_content is None and raw_part is None:
        await message.reply(
            f"Не вмію читати такий формат ({file_name}). "
            "Підтримую PDF, DOCX, XLSX, PPTX і звичайні текстові файли "
            "(txt, csv, json, md тощо)."
        )
        return

    question = strip_trigger(caption, BOT_USERNAME) or "Опрацюй цей файл і розкажи головне."
    question_text = f"{question}\n\n[Файл: {file_name}]"

    extra_parts = None
    if raw_part is not None:
        extra_parts = [raw_part]
    else:
        trimmed = text_content[:MAX_DOC_CHARS]
        if len(text_content) > MAX_DOC_CHARS:
            trimmed += "\n...[текст обрізано, файл завеликий]"
        question_text += f"\n\nВміст файлу:\n{trimmed}"

    await process_and_reply(
        message, sender, question_text,
        extra_parts=extra_parts,
        history_label=f"[файл {file_name}] {question}",
    )


def _validate_config() -> None:
    """Падаємо одразу зі зрозумілим повідомленням, якщо чогось не вистачає
    в .env — краще явна помилка при старті, ніж незрозумілий збій пізніше."""
    missing = []
    if not BOT_TOKEN:
        missing.append("BOT_TOKEN")
    if not GEMINI_API_KEY:
        missing.append("GEMINI_API_KEY")
    if not GEMINI_MODEL:
        missing.append("GEMINI_MODEL")
    if missing:
        raise RuntimeError(
            "Не задані обов'язкові змінні оточення: " + ", ".join(missing)
        )


async def main():
    global BOT_ID, BOT_USERNAME, BOT_FULL_NAME

    _validate_config()

    print(f"Поточна модель в боті: {GEMINI_MODEL}")
    log.info("Бот запускається...")

    me = await bot.get_me()
    BOT_ID = me.id
    BOT_USERNAME = me.username
    BOT_FULL_NAME = me.full_name
    log.info(f"Бот авторизований як @{BOT_USERNAME} (id={BOT_ID})")

    # Koyeb (безкоштовний план) вимагає Web Service з відкритим портом —
    # піднімаємо мінімальний HTTP-сервер для health-check поруч з polling'ом.
    port = int(os.getenv("PORT", "8000"))

    async def health(request):
        return web.Response(text="OK")

    app = web.Application()
    app.router.add_get("/", health)
    app.router.add_get("/health", health)

    runner = web.AppRunner(app)
    await runner.setup()
    site = web.TCPSite(runner, host="0.0.0.0", port=port)
    await site.start()
    log.info(f"Health-check сервер запущено на порту {port}")

    asyncio.create_task(idle_chat_watcher())
    log.info(f"Спостерігач за тишею в чаті запущено (поріг {IDLE_HOURS}г)")

    await dp.start_polling(bot)


if __name__ == "__main__":
    asyncio.run(main())
